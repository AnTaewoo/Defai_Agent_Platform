"""[P1] PowerPoint(.pptx) 파서. python-pptx로 직접 제어 — LangChain 로더 위임 금지."""

from __future__ import annotations

from pptx import Presentation

from ..types import ParsedDocument


def parse_pptx(path: str) -> ParsedDocument:
    """pptx 파일을 ParsedDocument로 변환.

    슬라이드마다 모든 shape의 텍스트 프레임 텍스트를 모아 슬라이드 1개당
    text_blocks 항목 하나로 담는다(슬라이드 번호를 맥락 정보로 앞에 붙임).
    슬라이드 내 표는 tables에 행렬(list[list[str]])로 보존하고, 빈 셀은 ""로 채운다.
    """
    prs = Presentation(path)

    text_blocks: list[str] = []
    tables: list[list[list[str]]] = []

    for idx, slide in enumerate(prs.slides, start=1):
        lines: list[str] = []

        for shape in slide.shapes:
            if shape.has_table:
                rows: list[list[str]] = []
                for row in shape.table.rows:
                    cells = ["" if cell.text is None else cell.text for cell in row.cells]
                    rows.append(cells)
                if rows:
                    tables.append(rows)
                continue

            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    lines.append(text)

        if lines:
            text_blocks.append(f"# Slide {idx}\n" + "\n".join(lines))

    return ParsedDocument(
        source=path,
        doc_type="pptx",
        text_blocks=text_blocks,
        tables=tables,
    )
