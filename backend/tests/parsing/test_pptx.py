from pptx import Presentation
from pptx.util import Inches

from src.parsing import parse
from src.parsing.pptx import parse_pptx


def _make_pptx(path: str) -> str:
    prs = Presentation()

    title_slide_layout = prs.slide_layouts[5]  # 제목만 있는 레이아웃

    # Slide 1: 텍스트
    slide1 = prs.slides.add_slide(title_slide_layout)
    slide1.shapes.title.text = "Hello D.A.P"

    # Slide 2: 표
    slide2 = prs.slides.add_slide(title_slide_layout)
    slide2.shapes.title.text = "표 슬라이드"
    rows, cols = 2, 3
    table_shape = slide2.shapes.add_table(
        rows, cols, Inches(1), Inches(1), Inches(4), Inches(2)
    )
    table = table_shape.table
    table.cell(0, 0).text = "이름"
    table.cell(0, 1).text = "부서"
    table.cell(0, 2).text = "등급"
    table.cell(1, 0).text = "홍길동"
    table.cell(1, 1).text = "기획"
    table.cell(1, 2).text = "2"

    prs.save(path)
    return path


def test_parse_pptx_extracts_text_and_table(tmp_path):
    path = _make_pptx(str(tmp_path / "sample.pptx"))

    doc = parse_pptx(path)

    assert doc.doc_type == "pptx"
    assert doc.source == path
    assert any("Hello D.A.P" in block for block in doc.text_blocks)
    assert any("표 슬라이드" in block for block in doc.text_blocks)
    assert doc.tables == [
        [["이름", "부서", "등급"], ["홍길동", "기획", "2"]]
    ]


def test_parse_dispatches_pptx_by_extension(tmp_path):
    path = _make_pptx(str(tmp_path / "sample.pptx"))

    doc = parse(path)

    assert doc.doc_type == "pptx"
